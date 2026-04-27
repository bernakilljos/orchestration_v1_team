#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate-final-ppt.py - HTML/CSS 슬라이드를 Playwright 로 PNG 렌더링 후
                        python-pptx 에 풀블리드로 조립.

Architecture:
  1. Playwright async: 1920×1080 viewport, device_scale_factor=2 (4K)
  2. 각 HTML → PNG 렌더링 (font load 대기)
  3. python-pptx: blank slide layout + 이미지 풀블리드 삽입

Usage:
  python generate-final-ppt.py [--dry-run] [--verbose]

Output:
  outputs/ppt/orchestration-v1-FINAL.pptx (25 slides)
"""

import asyncio
import sys
from pathlib import Path
from typing import List, Optional

try:
    from playwright.async_api import async_playwright, TimeoutError as PlaywrightTimeout
except ImportError:
    print("ERROR: playwright not installed. Run: pip install playwright")
    print("       Then: playwright install chromium")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ============================================================================
# CONFIGURATION
# ============================================================================

ROOT = Path(__file__).resolve().parent.parent.parent
HTML_DIR = ROOT / "outputs/ppt/html-source/slides"
PNG_DIR = ROOT / "outputs/ppt/html-source/png-output"
OUTPUT_PPTX = ROOT / "outputs/ppt/orchestration-v1-FINAL.pptx"

SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080
DEVICE_SCALE = 2  # 4K rendering for retina quality

DRY_RUN = "--dry-run" in sys.argv
VERBOSE = "--verbose" in sys.argv or "-v" in sys.argv

# ============================================================================
# UTILITIES
# ============================================================================

def log(msg: str, level: str = "INFO"):
    """Simple logging."""
    prefix = "OK" if level == "OK" else "!" if level == "WARN" else "->"
    try:
        print(f"[{prefix}] {msg}", flush=True)
    except (UnicodeEncodeError, TypeError):
        # Fallback for terminal encoding issues
        safe_msg = msg.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
        print(f"[{prefix}] {safe_msg}", flush=True)

def debug(msg: str):
    """Debug logging (verbose only)."""
    if VERBOSE:
        try:
            print(f"    {msg}", flush=True)
        except (UnicodeEncodeError, TypeError):
            safe_msg = msg.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
            print(f"    {safe_msg}", flush=True)

# ============================================================================
# PLAYWRIGHT RENDERING
# ============================================================================

async def render_all() -> List[Path]:
    """
    Render all HTML slides to PNG.

    Returns:
        List of PNG file paths in order.
    """
    PNG_DIR.mkdir(parents=True, exist_ok=True)

    slides = sorted(HTML_DIR.glob("slide-*.html"))
    if not slides:
        log(f"ERROR: No slides found in {HTML_DIR}", "WARN")
        return []

    log(f"Found {len(slides)} slides to render")

    if DRY_RUN:
        log(f"DRY-RUN: Would render {len(slides)} slides", "WARN")
        return [PNG_DIR / f"slide-{i:02d}.png" for i in range(1, len(slides) + 1)]

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        debug(f"Chromium launched")

        context = await browser.new_context(
            viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT},
            device_scale_factor=DEVICE_SCALE,
            ignore_https_errors=True,
        )
        page = await context.new_page()

        rendered_pngs = []

        for idx, slide_html in enumerate(slides, 1):
            try:
                log(f"[{idx}/{len(slides)}] Rendering {slide_html.name}...")

                # Navigate to HTML file
                file_url = slide_html.absolute().as_uri()
                await page.goto(file_url, wait_until="networkidle", timeout=10000)
                debug(f"  Page loaded: {file_url}")

                # Wait for fonts and images to load
                try:
                    await page.wait_for_function(
                        "document.fonts.ready",
                        timeout=5000
                    )
                    debug("  Fonts ready")
                except PlaywrightTimeout:
                    debug("  Font load timeout (continuing anyway)")

                # Wait for layout stability
                await asyncio.sleep(0.5)

                # Take screenshot
                png_path = PNG_DIR / f"slide-{idx:02d}.png"
                await page.screenshot(
                    path=str(png_path),
                    full_page=False,
                    type="png",
                    clip={
                        "x": 0,
                        "y": 0,
                        "width": SLIDE_WIDTH,
                        "height": SLIDE_HEIGHT,
                    },
                )
                rendered_pngs.append(png_path)
                log(f"  [OK] {png_path.name} ({png_path.stat().st_size / 1024:.1f} KB)")

            except Exception as e:
                log(f"  ERROR rendering {slide_html.name}: {e}", "WARN")
                continue

        await browser.close()
        debug("Browser closed")

    return rendered_pngs

# ============================================================================
# PPTX ASSEMBLY
# ============================================================================

def assemble_pptx(png_files: List[Path]) -> Path:
    """
    Assemble PNG files into PPTX with full-bleed layout.

    Args:
        png_files: List of PNG paths in order

    Returns:
        Output PPTX path
    """
    if not png_files:
        log("ERROR: No PNG files to assemble", "WARN")
        return None

    log(f"Creating PPTX with {len(png_files)} slides...")

    # Create presentation
    prs = Presentation()

    # Set slide dimensions to 1920×1080 (landscape widescreen)
    # EMU = English Metric Units. 1 inch = 914400 EMU
    # 1920px ≈ 2.67 inches, 1080px ≈ 1.5 inches (at 72 DPI)
    # But we want exact 1920×1080, so use EMU directly:
    # 1920 * 9525 = 18,288,000 EMU
    # 1080 * 9525 = 10,287,000 EMU
    prs.slide_width = Emu(int(SLIDE_WIDTH * 9525))
    prs.slide_height = Emu(int(SLIDE_HEIGHT * 9525))

    debug(f"Slide dimensions: {prs.slide_width} × {prs.slide_height} EMU")

    # Use blank layout (index 6)
    blank_layout = prs.slide_layouts[6]

    for png_file in png_files:
        if not png_file.exists():
            log(f"  WARN: {png_file.name} not found, skipping", "WARN")
            continue

        slide = prs.slides.add_slide(blank_layout)
        debug(f"  Added blank slide")

        # Insert image full-bleed
        slide.shapes.add_picture(
            str(png_file),
            Emu(0),
            Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        debug(f"  Inserted {png_file.name}")

    # Save
    prs.save(str(OUTPUT_PPTX))
    log(f"[OK] PPTX saved: {OUTPUT_PPTX}")

    return OUTPUT_PPTX

# ============================================================================
# MAIN
# ============================================================================

async def main():
    """Main entry point."""
    log("=" * 70)
    log("Orchestration Kit v1 - Final PPT Generation")
    log("=" * 70)
    log(f"HTML Source: {HTML_DIR}")
    log(f"PNG Output:  {PNG_DIR}")
    log(f"PPTX Output: {OUTPUT_PPTX}")
    if DRY_RUN:
        log("MODE: DRY-RUN (no files created)", "WARN")
    log("=" * 70)

    # Step 1: Render HTML → PNG
    log("\n[STEP 1] Rendering HTML to PNG...")
    png_files = await render_all()

    if not png_files:
        log("No PNGs generated. Aborting.", "WARN")
        return 1

    log(f"\n[OK] Rendered {len(png_files)} slides")

    if DRY_RUN:
        log("DRY-RUN: Skipping PPTX assembly", "WARN")
        return 0

    # Step 2: Assemble PPTX
    log("\n[STEP 2] Assembling PPTX...")
    output = assemble_pptx(png_files)

    if not output:
        log("Failed to create PPTX", "WARN")
        return 1

    # Verification
    log("\n[VERIFICATION]")
    if OUTPUT_PPTX.exists():
        size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)
        log(f"File exists: {OUTPUT_PPTX}")
        log(f"File size:   {size_mb:.2f} MB")
        log(f"Slides:      {len(png_files)}")
        log("[OK] PPT generation complete!")
        return 0
    else:
        log("ERROR: PPTX file not created", "WARN")
        return 1

if __name__ == "__main__":
    exit_code = asyncio.run(main())
    sys.exit(exit_code)
