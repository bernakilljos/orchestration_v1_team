#!/usr/bin/env python3
"""
render-ppt-multi.py — 3개 PPT (automation/plugins/team) HTML→PNG→PPTX.

Usage:
  python .claude/scripts/render-ppt-multi.py auto plugins team    # all 3
  python .claude/scripts/render-ppt-multi.py auto                 # only one
  python .claude/scripts/render-ppt-multi.py auto --slides 07 08 09  # subset of slides
"""
import asyncio
import sys
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent.parent

DECKS = {
    "auto": {
        "html": ROOT / "outputs/ppt-automation/html-source/slides",
        "png": ROOT / "outputs/ppt-automation/html-source/png-output",
        "pptx": ROOT / "outputs/ppt-automation/automation-guide.pptx",
    },
    "plugins": {
        "html": ROOT / "outputs/ppt-plugins/html-source/slides",
        "png": ROOT / "outputs/ppt-plugins/html-source/png-output",
        "pptx": ROOT / "outputs/ppt-plugins/plugins-guide.pptx",
    },
    "team": {
        "html": ROOT / "outputs/ppt-team/html-source/slides",
        "png": ROOT / "outputs/ppt-team/html-source/png-output",
        "pptx": ROOT / "outputs/ppt-team/team-guide.pptx",
    },
}

SLIDE_W, SLIDE_H, SCALE = 1920, 1080, 2


async def render_html(html_dir: Path, png_dir: Path, only: list[str] | None):
    png_dir.mkdir(parents=True, exist_ok=True)
    htmls = sorted(html_dir.glob("slide-*.html"))
    if only:
        htmls = [h for h in htmls if any(f"slide-{n}." in h.name or f"slide-{n}.html" == h.name for n in only)]
    print(f"  Rendering {len(htmls)} slide(s)...")

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": SLIDE_W, "height": SLIDE_H},
            device_scale_factor=SCALE,
        )
        page = await ctx.new_page()
        for html in htmls:
            url = html.resolve().as_uri()
            png = png_dir / f"{html.stem}.png"
            await page.goto(url, wait_until="networkidle", timeout=20000)
            try:
                await page.wait_for_load_state("networkidle", timeout=5000)
            except Exception:
                pass
            await page.wait_for_timeout(800)
            await page.screenshot(path=str(png), full_page=False, omit_background=False,
                                  clip={"x": 0, "y": 0, "width": SLIDE_W, "height": SLIDE_H})
            print(f"    {html.name} -> {png.name}")
        await browser.close()


def build_pptx(png_dir: Path, pptx_path: Path):
    pngs = sorted(png_dir.glob("slide-*.png"))
    print(f"  Assembling {len(pngs)} png(s) -> {pptx_path.name}")
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W * 9525)
    prs.slide_height = Emu(SLIDE_H * 9525)
    blank = prs.slide_layouts[6]
    for png in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))


async def main():
    args = sys.argv[1:]
    only_slides = None
    if "--slides" in args:
        i = args.index("--slides")
        only_slides = args[i+1:]
        args = args[:i]
    if not args:
        args = list(DECKS)

    for key in args:
        if key not in DECKS:
            print(f"SKIP unknown deck: {key}")
            continue
        d = DECKS[key]
        print(f"=== {key} ===")
        await render_html(d["html"], d["png"], only_slides)
        build_pptx(d["png"], d["pptx"])

    print("DONE")


if __name__ == "__main__":
    asyncio.run(main())
